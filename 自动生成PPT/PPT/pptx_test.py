#!/usr/bin/env python
# coding: utf-8

# In[1]:


import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_height = 6858000 
prs.slide_width = 12192000 
blank_slide_layout = prs.slide_layouts[6]
bullet_slide_layout = prs.slide_layouts[1]

def SetPictureAndTitle(slide,img_path):
    left, top, width, height = Inches(0), Inches(0), Inches(13.5), Inches(7.5)
    pic = slide.shapes.add_picture(img_path,left,top,width,height)
    left, top, width, height = Inches(3), Inches(0.1), Inches(8), Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.text = "主标题"
    p.font.size = Pt(40)

def start1(): 
    slide = prs.slides.add_slide(blank_slide_layout)
    #设置背景图片和主标题
    img_path = 'backgrounds/1_1.jpg'
    left, top, width, height = Inches(0), Inches(0), Inches(13.5), Inches(7.5)
    pic = slide.shapes.add_picture(img_path,left,top,width,height)
    #主标题
    left, top, width, height = Inches(1.5), Inches(4.5), Inches(8), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.text = "数据助力水处理"
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.size = Pt(60)
    #副标题1
    left = Inches(3.2)
    top = Inches(6)
    width = Inches(3)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.color.rgb = RGBColor(255,255,255)
    p.text = "XXX科技有限公司"
    p.font.size = Pt(25)
    #副标题2
    left = Inches(3.5)
    top = Inches(6.5)
    width = Inches(2)
    height = Inches(0.7)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.color.rgb = RGBColor(255,255,255)
    p.text = "日期：XXX年xxx月"
    p.font.size = Pt(15)
    prs.save('test.pptx')

def start2():
    
    slide = prs.slides.add_slide(blank_slide_layout)
    #设置背景图片和主标题
    img_path = 'backgrounds/2_1.jpg'
    left, top, width, height = Inches(0), Inches(0), Inches(13.5), Inches(7.5)
    pic = slide.shapes.add_picture(img_path,left,top,width,height)
    #主标题
    left, top, width, height = Inches(3), Inches(2), Inches(8), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.text = "数据助力水处理设施运维服务"
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.size = Pt(40)
    #副标题1
    left = Inches(5.5)
    top = Inches(5)
    width = Inches(3)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.text = "XXX科技有限公司"
    p.font.size = Pt(25)
    #副标题2
    left = Inches(5.8)
    top = Inches(5.8)
    width = Inches(3)
    height = Inches(0.7)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.text = "日期：XXX年xxx月"
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.size = Pt(15)
    prs.save('test.pptx')

def start3():   
    slide = prs.slides.add_slide(blank_slide_layout)
    #设置背景图片和主标题
    img_path = 'backgrounds/3_1.jpg'
    left, top, width, height = Inches(0), Inches(0), Inches(13.5), Inches(7.5)
    pic = slide.shapes.add_picture(img_path,left,top,width,height)
    #主标题
    left, top, width, height = Inches(5), Inches(2.5), Inches(5), Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.bold = True
    p.text = "工作总结"
    p.font.color.rgb = RGBColor(255,255,0)
    p.font.size = Pt(80)
    #副标题1
    left = Inches(7)
    top = Inches(5)
    width = Inches(3)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.color.rgb = RGBColor(255,255,0)
    p.text = "汇报人：XXXXX\n 日期：XXXXX"
    p.font.size = Pt(25)
    prs.save('test.pptx')

def directory1():
    slide = prs.slides.add_slide(blank_slide_layout)
    #设置背景图片和标题
    img_path = 'backgrounds/1_1.jpg'
    left, top, width, height = Inches(0), Inches(0), Inches(13.5), Inches(7.5)
    pic = slide.shapes.add_picture(img_path,left,top,width,height)
    #目录标题
    left, top, width, height = Inches(0.5), Inches(4), Inches(3), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.bold = True
    p.text = "目录"
    p.font.size = Pt(80)
    #目录
    left, top, width, height = Inches(3), Inches(3), Inches(3), Inches(3.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.text = "● 第一个目录\n"
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.text = "● 第二个目录\n"
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.text = "● 第三个目录\n"
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.text = "● 第四个目录\n"
    p = tf.add_paragraph()
    p.font.name = '黑体'
    prs.save('test.pptx')
    
def directory2():
    slide = prs.slides.add_slide(blank_slide_layout)
    #设置背景图片和标题
    img_path = 'backgrounds/2_2.jpg'
    left, top, width, height = Inches(0), Inches(0), Inches(13.5), Inches(7.5)
    pic = slide.shapes.add_picture(img_path,left,top,width,height)
    #目录标题
    left, top, width, height = Inches(1), Inches(0), Inches(3), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.font.bold = True
    p.text = "目录"
    p.font.size = Pt(50)
    
    #目录
    left, top, width, height = Inches(2.4), Inches(1.4), Inches(3), Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.size = Pt(20)
    p.text = "● 第一个目录\n"
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.size = Pt(20)
    p.text = "● 第二个目录\n"
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.size = Pt(20)
    p.text = "● 第三个目录\n"
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.size = Pt(20)
    p.text = "● 第四个目录\n"
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.size = Pt(20)
    p.text = "● 第五个目录\n"
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.size = Pt(20)
    p.text = "● 第六个目录\n"
    prs.save('test.pptx')
    
def directory3():
    slide = prs.slides.add_slide(blank_slide_layout)
    #设置背景图片和标题
    img_path = 'backgrounds/3_2.jpg'
    left, top, width, height = Inches(0), Inches(0), Inches(13.5), Inches(7.5)
    pic = slide.shapes.add_picture(img_path,left,top,width,height)
    
    #目录
    left, top, width, height = Inches(7), Inches(2.1), Inches(3), Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(208,181,91)
    p.text = "● 第一个目录\n"
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(208,181,91)
    p.text = "● 第二个目录\n"
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(208,181,91)
    p.text = "● 第三个目录\n"
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(208,181,91)
    p.text = "● 第四个目录\n"
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(208,181,91)
    p.text = "● 第五个目录\n"
    prs.save('test.pptx')

def first_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #设置背景图片和标题
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/1.jpg'
    left, top, width, height = Inches(7), Inches(2.1), Inches(5), Inches(3)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落设置-读入数据
    txt_path = 'texts/1.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落设置-位置
    left, top, width, height = Inches(1), Inches(2), Inches(5), Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 1.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
    prs.save('test.pptx')

def second_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #设置背景图片和标题
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/2.jpg'
    left, top, width, height = Inches(1), Inches(2.1), Inches(4.2), Inches(3)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落设置-读入数据
    txt_path = 'texts/2.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落设置-位置
    left, top, width, height = Inches(5.5), Inches(2), Inches(6), Inches(6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 2.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
    prs.save('test.pptx')

def third_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'images/3.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/3_1.jpg'
    left, top, width, height = Inches(10), Inches(5.6), Inches(2.3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落设置-读入数据
    txt_path = 'texts/3.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落设置-位置
    left, top, width, height = Inches(0.5), Inches(2), Inches(10), Inches(6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 3.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.color.rgb = RGBColor(255,255,255)
        p.font.name = '黑体'
        p.text = line
    prs.save('test.pptx')

def fouth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/4.jpg'
    left, top, width, height = Inches(0.5), Inches(1.7), Inches(4), Inches(5.5)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落设置-读入数据
    txt_path = 'texts/4.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落设置-位置
    left, top, width, height = Inches(4.5), Inches(1.5), Inches(6), Inches(6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 4.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
    prs.save('test.pptx')

def fifth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/5.jpg'
    left, top, width, height = Inches(7), Inches(2.5), Inches(5), Inches(3)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落1设置-读入数据
    txt_path = 'texts/5.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(0.2), Inches(2), Inches(6), Inches(6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 5.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
    left, top, width, height = Inches(7), Inches(6), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "这里可以输入文字"
    prs.save('test.pptx')

def sixth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/6.jpg'
    left, top, width, height = Inches(9), Inches(1.5), Inches(2), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #图标1+文字位置
    icon_path = 'icons/1.jpg'
    left, top, width, height = Inches(2), Inches(3.5), Inches(0.7), Inches(0.7)
    pic = slide.shapes.add_picture(icon_path, left, top, width, height)
    
    left, top, width, height = Inches(2.8), Inches(3.5), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第一个标题"
    #图标2+文字位置
    icon_path = 'icons/2.jpg'
    left, top, width, height = Inches(2), Inches(5.3), Inches(0.7), Inches(0.7)
    pic = slide.shapes.add_picture(icon_path, left, top, width, height)
    
    left, top, width, height = Inches(2.8), Inches(5.3), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第二个标题"
    #图标3+文字位置
    icon_path = 'icons/3.jpg'
    left, top, width, height = Inches(6.5), Inches(3.5), Inches(0.7), Inches(0.7)
    pic = slide.shapes.add_picture(icon_path, left, top, width, height)
    
    left, top, width, height = Inches(7.3), Inches(3.5), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第三个标题"
    #图标4+文字位置
    icon_path = 'icons/4.jpg'
    left, top, width, height = Inches(6.5), Inches(5.3), Inches(0.7), Inches(0.7)
    pic = slide.shapes.add_picture(icon_path, left, top, width, height)
    
    left, top, width, height = Inches(7.3), Inches(5.3), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第四个标题"
    prs.save('test.pptx')

def seventh_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片1位置
    img_path = 'images/7_1.jpg'
    left, top, width, height = Inches(2), Inches(4.4), Inches(2), Inches(3)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #图片2位置
    img_path = 'images/7_2.jpg'
    left, top, width, height = Inches(4.5), Inches(4.4), Inches(2), Inches(3)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #图片3位置
    img_path = 'images/7_3.jpg'
    left, top, width, height = Inches(7), Inches(4.4), Inches(2), Inches(3)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落1设置-读入数据
    txt_path = 'texts/7.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(2), Inches(2), Inches(8), Inches(2.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 7.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
    prs.save('test.pptx')

def eighth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片1位置
    img_path = 'images/8_1.jpg'
    left, top, width, height = Inches(1.3), Inches(2), Inches(3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #图片2位置
    img_path = 'images/8_2.jpg'
    left, top, width, height = Inches(5), Inches(2), Inches(3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #图片3位置
    img_path = 'images/8_3.jpg'
    left, top, width, height = Inches(8.7), Inches(2), Inches(3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落1设置-读入数据
    txt_path = 'texts/8.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(1.3), Inches(4), Inches(9.2), Inches(2.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 8.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
        p.font.size = Pt(15)
    prs.save('test.pptx')

def ninth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片1位置
    img_path = 'images/9.jpg'
    left, top, width, height = Inches(2), Inches(2), Inches(4), Inches(3)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落1设置-读入数据
    txt_path = 'texts/9_1.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(6.2), Inches(2), Inches(6), Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 9.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
    #段落2设置-读入数据
    txt_path = 'texts/9_2.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落2设置-位置
    left, top, width, height = Inches(2), Inches(5), Inches(4), Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
        p.font.size = Pt(25)
    prs.save('test.pptx')

def tenth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/10.jpg'
    left, top, width, height = Inches(0), Inches(2), Inches(13.5), Inches(6.5)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落1设置-读入数据
    txt_path = 'texts/10.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(2), Inches(2), Inches(4), Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 10.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
        p.font.size = Pt(25)
        p.font.bold = True
    prs.save('test.pptx')

def eleventh_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/11.jpg'
    left, top, width, height = Inches(7), Inches(2), Inches(5), Inches(3.5)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落1设置-读入数据
    txt_path = 'texts/11_1.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(0.3), Inches(2), Inches(5.5), Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 11_1和11_2.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
        p.font.size = Pt(15)
    #段落2设置-读入数据
    txt_path = 'texts/11_2.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落2设置-位置
    left, top, width, height = Inches(0.3), Inches(5.8), Inches(11), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
        p.font.size = Pt(18)
    prs.save('test.pptx')

def twelfth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片1位置
    img_path = 'images/12_1.jpg'
    left, top, width, height = Inches(0.2), Inches(3.5), Inches(3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    
    left, top, width, height = Inches(1), Inches(2), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第一个标题"
    
    #图片2位置
    img_path = 'images/12_2.jpg'
    left, top, width, height = Inches(3.4), Inches(3.5), Inches(3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    
    left, top, width, height = Inches(4.2), Inches(2), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第二个标题"
    
    #图片3位置
    img_path = 'images/12_3.jpg'
    left, top, width, height = Inches(6.6), Inches(3.5), Inches(3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    
    left, top, width, height = Inches(7.4), Inches(2), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第三个标题"
    
    #图片4位置
    img_path = 'images/12_4.jpg'
    left, top, width, height = Inches(9.8), Inches(3.5), Inches(3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    
    left, top, width, height = Inches(10.6), Inches(2), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第二个标题"
    prs.save('test.pptx')
    
def thirteenth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #段落1设置-读入数据
    txt_path = 'texts/13.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(0.5), Inches(2), Inches(12), Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 13.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
        p.font.size = Pt(15)
    prs.save('test.pptx')

def fouteenth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #设置背景图片和标题
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/1.jpg'
    left, top, width, height = Inches(7), Inches(2.1), Inches(5), Inches(3)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落设置-读入数据
    txt_path = 'texts/14.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落设置-位置
    left, top, width, height = Inches(1), Inches(2), Inches(5), Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 14.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
    prs.save('test.pptx')

def fifteenth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #设置背景图片和标题
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/15.jpg'
    left, top, width, height = Inches(1), Inches(2.1), Inches(4.2), Inches(3)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落设置-读入数据
    txt_path = 'texts/15.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落设置-位置
    left, top, width, height = Inches(5.5), Inches(2), Inches(6), Inches(6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 15.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
    prs.save('test.pptx')

def sixteenth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'images/3.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/16_1.jpg'
    left, top, width, height = Inches(10), Inches(5.6), Inches(2.3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落设置-读入数据
    txt_path = 'texts/16.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落设置-位置
    left, top, width, height = Inches(0.5), Inches(2), Inches(10), Inches(6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 16.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.color.rgb = RGBColor(255,255,255)
        p.font.name = '黑体'
        p.text = line
    prs.save('test.pptx')

def seventeeth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/17.jpg'
    left, top, width, height = Inches(0.5), Inches(1.7), Inches(4), Inches(5.5)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落设置-读入数据
    txt_path = 'texts/17.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落设置-位置
    left, top, width, height = Inches(4.5), Inches(1.5), Inches(6), Inches(6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 17.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
    prs.save('test.pptx')

def eighteenth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/18.jpg'
    left, top, width, height = Inches(7), Inches(2.5), Inches(5), Inches(3)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落1设置-读入数据
    txt_path = 'texts/18.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(0.2), Inches(2), Inches(6), Inches(6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 18.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
    left, top, width, height = Inches(7), Inches(6), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "这里可以输入文字"
    prs.save('test.pptx')

def nineteenth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/19.jpg'
    left, top, width, height = Inches(9), Inches(1.5), Inches(2), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #图标1+文字位置
    icon_path = 'icons/1.jpg'
    left, top, width, height = Inches(2), Inches(3.5), Inches(0.7), Inches(0.7)
    pic = slide.shapes.add_picture(icon_path, left, top, width, height)
    
    left, top, width, height = Inches(2.8), Inches(3.5), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第一个标题"
    #图标2+文字位置
    icon_path = 'icons/2.jpg'
    left, top, width, height = Inches(2), Inches(5.3), Inches(0.7), Inches(0.7)
    pic = slide.shapes.add_picture(icon_path, left, top, width, height)
    
    left, top, width, height = Inches(2.8), Inches(5.3), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第二个标题"
    #图标3+文字位置
    icon_path = 'icons/3.jpg'
    left, top, width, height = Inches(6.5), Inches(3.5), Inches(0.7), Inches(0.7)
    pic = slide.shapes.add_picture(icon_path, left, top, width, height)
    
    left, top, width, height = Inches(7.3), Inches(3.5), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第三个标题"
    #图标4+文字位置
    icon_path = 'icons/4.jpg'
    left, top, width, height = Inches(6.5), Inches(5.3), Inches(0.7), Inches(0.7)
    pic = slide.shapes.add_picture(icon_path, left, top, width, height)
    
    left, top, width, height = Inches(7.3), Inches(5.3), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第四个标题"
    prs.save('test.pptx')

def twentieth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片1位置
    img_path = 'images/20_1.jpg'
    left, top, width, height = Inches(2), Inches(4.4), Inches(2), Inches(3)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #图片2位置
    img_path = 'images/20_2.jpg'
    left, top, width, height = Inches(4.5), Inches(4.4), Inches(2), Inches(3)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #图片3位置
    img_path = 'images/20_3.jpg'
    left, top, width, height = Inches(7), Inches(4.4), Inches(2), Inches(3)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落1设置-读入数据
    txt_path = 'texts/20.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(2), Inches(2), Inches(8), Inches(2.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 20.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
    prs.save('test.pptx')

def twenty_first_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片1位置
    img_path = 'images/21_1.jpg'
    left, top, width, height = Inches(1.3), Inches(2), Inches(3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #图片2位置
    img_path = 'images/21_2.jpg'
    left, top, width, height = Inches(5), Inches(2), Inches(3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #图片3位置
    img_path = 'images/21_3.jpg'
    left, top, width, height = Inches(8.7), Inches(2), Inches(3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落1设置-读入数据
    txt_path = 'texts/21.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(1.3), Inches(4), Inches(9.2), Inches(2.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 21.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
        p.font.size = Pt(15)
    prs.save('test.pptx')

def twenty_second_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片1位置
    img_path = 'images/22.jpg'
    left, top, width, height = Inches(2), Inches(2), Inches(4), Inches(3)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落1设置-读入数据
    txt_path = 'texts/22_1.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(6.2), Inches(2), Inches(6), Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 22_1和22_2f.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
    #段落2设置-读入数据
    txt_path = 'texts/22_2.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落2设置-位置
    left, top, width, height = Inches(2), Inches(5), Inches(4), Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
        p.font.size = Pt(25)
    prs.save('test.pptx')

def twenty_third_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/23.jpg'
    left, top, width, height = Inches(0), Inches(2), Inches(13.5), Inches(6.5)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落1设置-读入数据
    txt_path = 'texts/23.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(2), Inches(2), Inches(4), Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 23.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
        p.font.size = Pt(25)
        p.font.bold = True
    prs.save('test.pptx')

def twenty_fouth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/24.jpg'
    left, top, width, height = Inches(7), Inches(2), Inches(5), Inches(3.5)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落1设置-读入数据
    txt_path = 'texts/24_1.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(0.3), Inches(2), Inches(5.5), Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 24.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
        p.font.size = Pt(15)
    #段落2设置-读入数据
    txt_path = 'texts/24_2.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落2设置-位置
    left, top, width, height = Inches(0.3), Inches(5.8), Inches(11), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
        p.font.size = Pt(18)
    prs.save('test.pptx')

def twenty_fifth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片1位置
    img_path = 'images/25_1.jpg'
    left, top, width, height = Inches(0.2), Inches(3.5), Inches(3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    
    left, top, width, height = Inches(1), Inches(2), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第一个标题"
    
    #图片2位置
    img_path = 'images/25_2.jpg'
    left, top, width, height = Inches(3.4), Inches(3.5), Inches(3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    
    left, top, width, height = Inches(4.2), Inches(2), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第二个标题"
    
    #图片3位置
    img_path = 'images/25_3.jpg'
    left, top, width, height = Inches(6.6), Inches(3.5), Inches(3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    
    left, top, width, height = Inches(7.4), Inches(2), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第三个标题"
    
    #图片4位置
    img_path = 'images/25_4.jpg'
    left, top, width, height = Inches(9.8), Inches(3.5), Inches(3), Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    
    left, top, width, height = Inches(10.6), Inches(2), Inches(2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.add_paragraph()
    p.font.name = '黑体'
    p.text = "第二个标题"
    prs.save('test.pptx')
    
def twenty_sixth_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #段落1设置-读入数据
    txt_path = 'texts/26.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(0.5), Inches(2), Inches(12), Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 26.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
        p.font.size = Pt(15)
    prs.save('test.pptx')
    
def twenty_seventh_page():
    slide = prs.slides.add_slide(blank_slide_layout)
    #背景图片
    img_path = 'backgrounds/1.jpg'
    SetPictureAndTitle(slide,img_path)
    #图片位置
    img_path = 'images/27.jpg'
    left, top, width, height = Inches(7), Inches(2), Inches(5), Inches(3.5)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    #段落1设置-读入数据
    txt_path = 'texts/27_1.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落1设置-位置
    left, top, width, height = Inches(0.3), Inches(2), Inches(5.5), Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "这里控制 27_1 和 27_2.txt"
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
        p.font.size = Pt(15)
    #段落2设置-读入数据
    txt_path = 'texts/27_2.txt'
    text = open(txt_path,encoding = 'utf-8')
    lines = text.readlines()
    #段落2设置-位置
    left, top, width, height = Inches(0.3), Inches(5.8), Inches(11), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in lines:
        p = tf.add_paragraph()
        p.font.name = '黑体'
        p.text = line
        p.font.size = Pt(18)
    prs.save('test.pptx')

def end1():
    slide = prs.slides.add_slide(blank_slide_layout)
    #设置背景图片和主标题
    img_path = 'backgrounds/1_3.jpg'
    left, top, width, height = Inches(0), Inches(0), Inches(13.5), Inches(7.5)
    pic = slide.shapes.add_picture(img_path,left,top,width,height)
    #主标题
    left, top, width, height = Inches(1), Inches(2.2), Inches(8), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.text = "谢谢聆听"
    p.font.size = Pt(40)
    prs.save('test.pptx')

def end2():
    slide = prs.slides.add_slide(blank_slide_layout)
    #设置背景图片和主标题
    img_path = 'backgrounds/2_3.jpg'
    left, top, width, height = Inches(0), Inches(0), Inches(13.5), Inches(7.5)
    pic = slide.shapes.add_picture(img_path,left,top,width,height)
    #主标题
    left, top, width, height = Inches(5), Inches(1.5), Inches(8), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.text = "谢谢聆听"
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.size = Pt(80)
    #副标题1
    left = Inches(5.5)
    top = Inches(5)
    width = Inches(3)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.text = "XXX科技有限公司"
    p.font.size = Pt(25)
    #副标题2
    left = Inches(5.8)
    top = Inches(5.8)
    width = Inches(3)
    height = Inches(0.7)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.text = "日期：XXX年xxx月"
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.size = Pt(15)
    prs.save('test.pptx')
    
def end3():
    slide = prs.slides.add_slide(blank_slide_layout)
    #设置背景图片和主标题
    img_path = 'backgrounds/3_3.jpg'
    left, top, width, height = Inches(0), Inches(0), Inches(13.5), Inches(7.5)
    pic = slide.shapes.add_picture(img_path,left,top,width,height)
    #主标题
    left, top, width, height = Inches(4), Inches(4.6), Inches(4), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.name = '黑体'
    p.text = "谢谢聆听"
    p.font.size = Pt(80)
    prs.save('test.pptx')

start1()
start2()
start3()
directory1()
directory2()
directory3()

def RandomExcute():
    function = {
    1:"first_page()",
    2:"second_page()",
    3:"third_page()",
    4:"fouth_page()",
    5:"fifth_page()",
    6:"sixth_page()",
    7:"seventh_page()",
    8:"eighth_page()",
    9:"ninth_page()",
    10:"tenth_page()",
    11:"eleventh_page()",
    12:"twelfth_page()",
    13:"thirteenth_page()"
    }
    randomList = random.sample(range(1,14),13)
    for r in randomList:
        exec(function[r])
        
RandomExcute()

'''
first_page()
second_page()
third_page()
fouth_page()
fifth_page()
sixth_page()
seventh_page()
eighth_page()
ninth_page()
tenth_page()
eleventh_page()
twelfth_page()
thirteenth_page()
'''
fouteenth_page()
fifteenth_page()
sixteenth_page()
seventeeth_page()
eighteenth_page()
nineteenth_page()
twentieth_page()
twenty_fifth_page()
twenty_second_page()
twenty_third_page()
twenty_fouth_page()
twenty_fifth_page()
twenty_sixth_page()
twenty_seventh_page()




end1()
end2()
end3()